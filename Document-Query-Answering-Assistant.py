import gradio as gr
import spacy
import pdfplumber
import torch
import numpy as np
import faiss
import mimetypes
import pytesseract
from pdf2image import convert_from_path
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer, AutoModelForQuestionAnswering, AutoModelForCausalLM
from sentence_transformers import SentenceTransformer
import pandas as pd
from docx import Document
from pptx import Presentation
import re
import cv2
import os

# Load NLP models
nlp = spacy.load("en_core_web_md")
device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
embedder = SentenceTransformer('all-MiniLM-L6-v2', device=device)

model_name = "google/flan-t5-base"  # Generative model for reasoning
tokenizer = AutoTokenizer.from_pretrained(model_name)

model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
model.to(device)

# FAISS Index Storage
faiss_index = None
stored_sentences = []
stored_embeddings = None

def check_valid_document(file):
    """Check if file type is supported."""
    if file is None:
        return False
    mime_type, _ = mimetypes.guess_type(file.name)
    allowed_types = [
        "application/pdf", "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "text/plain", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "text/csv",
        "application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ]
    return mime_type in allowed_types

def preprocess_image_for_ocr(image):
    """Preprocess images for better OCR performance."""
    gray = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)
    _, binary = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return binary

def extract_text_from_file(file):
    """Extract text from PDFs, Word, CSV, or Excel files efficiently."""
    text = ""
    mime_type, _ = mimetypes.guess_type(file.name)

    try:
        if mime_type == "application/pdf":
            with pdfplumber.open(file) as pdf:
                for page in pdf.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += re.sub(r'\s+', ' ', extracted).strip() + "\n\n"

            # **Only run OCR if text is missing**
            if not text.strip():
                images = convert_from_path(file.name, first_page=1, last_page=2)  # Limit to first 2 pages
                if images:
                    text = "\n\n".join([pytesseract.image_to_string(preprocess_image_for_ocr(img)) for img in images])

        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            text = "\n\n".join([para.text.strip() for para in doc.paragraphs])

        elif mime_type == "text/plain":
            text = file.read().decode("utf-8").strip()

        elif mime_type in ["application/vnd.ms-powerpoint",
                           "application/vnd.openxmlformats-officedocument.presentationml.slideshow",
                           "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            if mime_type == "application/vnd.ms-powerpoint":
                file = file.name.replace(".ppt", ".pptx")
            ppt = Presentation(file)
            text = "\n\n".join([slide.shapes.title.text for slide in ppt.slides])

        elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "text/csv"]:
            df = pd.read_excel(file) if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" else pd.read_csv(file)
            text = df.to_string(index=False)

    except Exception as e:
        return f"❌ Error extracting text: {e}"

    return text if text.strip() else "❌ No readable text found."

def build_faiss_index(text):
    """Build FAISS index using paragraph-based chunks for faster retrieval."""
    global faiss_index, stored_sentences, stored_embeddings

    if not text.strip():
        return "❌ No valid text."

    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    stored_sentences = paragraphs

    # Remove stored_embeddings caching (caused missing updates)
    stored_embeddings = np.array([embedder.encode(p, normalize_embeddings=True) for p in stored_sentences])

    if stored_embeddings.shape[0] == 0:
        return "❌ No valid embeddings."

    # Switch back to IndexFlatIP for better accuracy
    faiss_index = faiss.IndexFlatIP(stored_embeddings.shape[1])
    faiss_index.add(stored_embeddings)

    return "FAISS index built successfully!"

query_cache = {}

def retrieve_relevant_passages(query):
    """Retrieve relevant passages from FAISS index efficiently."""
    if faiss_index is None or not stored_sentences:
        return ["No indexed documents available."]

    query_embedding = embedder.encode([query], normalize_embeddings=True)

    # Lower threshold to 0.4 (previously 0.45) for better recall
    D, I = faiss_index.search(query_embedding, k=10)  # k=10 for more relevant results

    results = [stored_sentences[i] for i, score in zip(I[0], D[0]) if i >= 0 and score >= 0.30]

    return results if results else ["The document does not contain relevant information."]

def chunk_text(text, max_length=512, overlap=100):
    """Splits text into overlapping chunks for long context handling."""
    words = text.split()
    chunks = []
    start = 0

    while start < len(words):
        chunk = words[start:start + max_length]
        chunks.append(" ".join(chunk))
        start += max_length - overlap  # Move forward with some overlap

    return chunks

def generate_answer(question):
    retrieved_passages = retrieve_relevant_passages(question)

    if not retrieved_passages or "The document does not contain relevant information." in retrieved_passages:
        return "The document does not provide an answer."

    context = " ".join(retrieved_passages)[:2048]  # Increased context length

    # Optimized Prompt
    prompt = f"""
You are an AI assistant that answers **strictly based on the document**.
If an **exact answer exists**, return it **verbatim**.
If an answer is **implied or partially available**, explain it using logical reasoning.

📌 **Answering Process:**
1️⃣ **Search for an Exact Match.**
   - If found, return the exact answer **verbatim**.

2️⃣ **Look for Reworded or Indirect Answers.**
   - If no direct answer exists, search for paraphrased or implied meanings.
   - If found, **explain how it relates to the user's question.**

3️⃣ **Scan Headers First, Then the Body.**
   - If the question relates to policies, rules, or frameworks, check **section titles**.
   - **If no direct match in headers, scan full sections under those headers**.

4️⃣ **Use Keyword Expansion for Better Matching.**
   - Look for **synonyms** (e.g., "frameworks" → "guidelines," "principles").
   - If the term is **not found exactly**, check **phrases with similar meanings**.

📌 **Final Rule:** If no exact answer exists:
   - **Infer the best possible answer from related content.**
   - If partial details exist, **explain the closest meaning**.
   - If no relevant content exists, return:
     📌 _"The document does not explicitly mention '{question}', but based on the closest details, a possible explanation is: {context}"_
     -If an answer is already generated, there's no need to provide more details.

Now summarize the key information from the retrieved sections.

---

### **User Question:**
{question}

### **Extracted Relevant Section(s):**
{context}

### **Final Answer:**
If no exact match exists, **explain the best possible meaning** based on related details.
If an answer is **implied, provide a logical explanation**.
If no details exist, **generate a possible reasoning based on the topic**.
"""

    # Model Generation (Ensures Answer Only)
    inputs = tokenizer(prompt + "Answer:\n", return_tensors="pt", truncation=True, max_length=1024).to(device)

    with torch.no_grad():
        outputs = model.generate(**inputs, max_length=200, num_beams=4, do_sample=False, repetition_penalty=1.8)

        if torch.cuda.is_available():
            torch.cuda.empty_cache()

    answer = tokenizer.decode(outputs[0], skip_special_tokens=True).strip()

    return answer if len(answer) > 5 else "The document does not provide an explicit answer, but here’s the closest information:\n\n" + context

# Gradio and interaction functions

import time

def chat_with_document(message, history):
    response = generate_answer(message)
    history.append((message, response))
    return history, history


def upload_and_process_file(file):
    """File processing and generation of document-based responses."""
    text = extract_text_from_file(file)
    build_faiss_index(text)

    return "Document uploaded and indexed successfully."

iface = gr.Interface(
    fn=chat_with_document,
    inputs=[gr.Textbox(placeholder="Ask a question about the document...")],
    outputs=[gr.Chatbot()],
    title="Document Query and Answering System",
    description="Upload a document and interact with it.",
    theme="compact"
)

iface.launch(debug=True)
