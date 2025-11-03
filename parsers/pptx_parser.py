"""PPTX document parser."""
from pathlib import Path
from typing import Dict
from pptx import Presentation
from .document_parser import DocumentParser

class PPTXParser(DocumentParser):
    """Parser for PPTX files."""
    
    def parse(self, file_path: Path) -> Dict[str, str]:
        """Extract text from PPTX file."""
        text = ""
        metadata = {"slides": 0}
        
        try:
            ppt = Presentation(file_path)
            metadata["slides"] = len(ppt.slides)
            
            for slide in ppt.slides:
                slide_text = []
                # Get title
                if slide.shapes.title and slide.shapes.title.text:
                    slide_text.append(f"Title: {slide.shapes.title.text}")
                
                # Get text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if shape.text.strip() and shape.text not in slide_text:
                            slide_text.append(shape.text.strip())
                
                if slide_text:
                    text += "\n".join(slide_text) + "\n\n"
            
            if not text.strip():
                text = "❌ No readable text found in PPTX."
        except Exception as e:
            text = f"❌ Error extracting text from PPTX: {str(e)}"
        
        return {"text": text.strip(), "metadata": metadata}

